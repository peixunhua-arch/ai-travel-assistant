"""途灵 · 精致视觉版会议汇报 PPT。"""
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# Brand tokens (对齐 App 主题)
CREAM = RGBColor(0xFF, 0xF8, 0xF0)
CREAM_DEEP = RGBColor(0xFF, 0xF0, 0xE4)
NAVY = RGBColor(0x1A, 0x2E, 0x32)
TEAL = RGBColor(0x3D, 0xB8, 0xA9)
TEAL_DARK = RGBColor(0x2E, 0x9A, 0x8E)
TEAL_DEEP = RGBColor(0x0C, 0x2E, 0x30)
INK = RGBColor(0x4F, 0x34, 0x28)
MUTED = RGBColor(0x96, 0x78, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xF5, 0xA0, 0x6A)
CARD = RGBColor(0xFF, 0xFC, 0xF9)
SOFT = RGBColor(0xE8, 0xF8, 0xF5)

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.65)
DOCS = Path(__file__).resolve().parent
SHOTS = DOCS / "screenshots"
TOTAL = 12


def font(run, size, bold=False, color=INK, name="Microsoft YaHei"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", name)


def bg(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp = shape._element
    slide.shapes._spTree.remove(sp)
    slide.shapes._spTree.insert(2, sp)


def rect(slide, l, t, w, h, color, radius=0.1):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def txt(slide, l, t, w, h, text, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    font(r, size, bold, color)
    return box


def footer(slide, page, light=False):
    c = RGBColor(0xC4, 0xB0, 0xA0) if not light else RGBColor(0x7A, 0x9A, 0x98)
    txt(slide, M, Inches(7.1), Inches(5), Inches(0.28), "途灵 · AI 旅游助手", 11, False, c)
    txt(slide, Inches(11.4), Inches(7.1), Inches(1.3), Inches(0.28), f"{page:02d} / {TOTAL:02d}", 11, False, c, PP_ALIGN.RIGHT)


def pic(slide, path: Path, l, t, height=None, width=None):
    if not path.exists():
        return None
    kw = {}
    if height is not None:
        kw["height"] = height
    if width is not None:
        kw["width"] = width
    return slide.shapes.add_picture(str(path), l, t, **kw)


def title_block(slide, title, subtitle=None, dark=False):
    ink = WHITE if dark else NAVY
    mute = RGBColor(0xA8, 0xC8, 0xC4) if dark else MUTED
    rect(slide, M, Inches(0.42), Inches(0.1), Inches(0.38), TEAL, 0.5)
    txt(slide, Inches(0.95), Inches(0.35), Inches(11), Inches(0.5), title, 26, True, ink)
    if subtitle:
        txt(slide, Inches(0.95), Inches(0.88), Inches(11), Inches(0.35), subtitle, 13, False, mute)


def pill(slide, l, t, w, h, text, fill=TEAL, fg=WHITE):
    rect(slide, l, t, w, h, fill, 0.5)
    txt(slide, l, t + Inches(0.08), w, h, text, 12, True, fg, PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    hq_guide = SHOTS / "hq-guide.png"
    hq_trips = SHOTS / "hq-trips.png"
    hq_chat = SHOTS / "hq-chat.png"
    hq_comm = SHOTS / "hq-community.png"
    hq_splash = SHOTS / "hq-splash.png"
    stage_splash = SHOTS / "stage-splash.png"
    poster = SHOTS / "poster-splash.png"
    collage = SHOTS / "hero-collage.png"
    icon = SHOTS / "icon.png"
    # fallbacks
    splash_phone = SHOTS / "00-splash-phone.png"
    if not hq_splash.exists():
        hq_splash = splash_phone
    if not hq_guide.exists():
        hq_guide = SHOTS / "01-plan.png"
    if not hq_trips.exists():
        hq_trips = SHOTS / "03-trips.png"
    if not hq_chat.exists():
        hq_chat = SHOTS / "06-chat.png"
    if not hq_comm.exists():
        hq_comm = SHOTS / "04-community.png"

    # ========== 1 Cover ==========
    s = prs.slides.add_slide(blank)
    bg(s, TEAL_DEEP)
    # left cream panel for brand
    rect(s, 0, 0, Inches(6.2), H, CREAM, 0)
    if icon.exists():
        pic(s, icon, Inches(0.85), Inches(1.55), height=Inches(0.7))
    txt(s, Inches(0.85), Inches(2.5), Inches(5), Inches(0.4), "项目汇报", 14, False, MUTED)
    txt(s, Inches(0.85), Inches(2.95), Inches(5), Inches(0.9), "途灵", 52, True, INK)
    txt(s, Inches(0.85), Inches(3.9), Inches(5), Inches(0.45), "你的 AI 出行助手", 18, False, TEAL_DARK)
    txt(
        s,
        Inches(0.85),
        Inches(4.55),
        Inches(4.8),
        Inches(0.8),
        "从规划到落地的完整链路\nAgent · 真实 POI · 评价闭环",
        14,
        False,
        MUTED,
    )
    pill(s, Inches(0.85), Inches(5.7), Inches(2.2), Inches(0.42), "可演示 · 可安装", TEAL, WHITE)
    # right: collage or phones
    if collage.exists():
        pic(s, collage, Inches(5.9), Inches(0.35), height=Inches(6.8))
    else:
        for i, p in enumerate([hq_splash, hq_guide, hq_trips]):
            if p.exists():
                pic(s, p, Inches(6.2 + i * 2.15), Inches(0.9), height=Inches(5.5))

    # ========== 2 Splash brand ==========
    s = prs.slides.add_slide(blank)
    bg(s, CREAM)
    title_block(s, "品牌开机页", "行李箱 + 定位 Pin · 第一眼建立「途灵」认知")
    points = [
        ("品牌锚点", "途灵 · 你的 AI 出行助手"),
        ("视觉语言", "暖奶油底 + 青绿路线，贯穿全 App"),
        ("产品隐喻", "行程可打包、路径可落地"),
        ("工程对应", "apps/mobile/assets/splash.png"),
    ]
    for i, (a, b) in enumerate(points):
        top = Inches(1.7) + Inches(i * 1.1)
        rect(s, M, top, Inches(6.4), Inches(0.95), WHITE, 0.12)
        txt(s, Inches(1.0), top + Inches(0.18), Inches(5.8), Inches(0.35), a, 14, True, TEAL_DARK)
        txt(s, Inches(1.0), top + Inches(0.5), Inches(5.8), Inches(0.35), b, 15, False, INK)
    phone = stage_splash if stage_splash.exists() else hq_splash
    if phone.exists():
        pic(s, phone, Inches(8.2), Inches(1.15), height=Inches(5.6))
    elif poster.exists():
        pic(s, poster, Inches(8.4), Inches(1.2), height=Inches(5.5))
    footer(s, 2)

    # ========== 3 Agenda ==========
    s = prs.slides.add_slide(blank)
    bg(s, CREAM)
    title_block(s, "汇报提纲", "约 6–8 分钟 · 界面与产品一体讲述")
    items = [
        ("01", "背景目标", "攻略散、假行程、难落地"),
        ("02", "产品方案", "Agent + 真数据 + App"),
        ("03", "界面演示", "高清水效果一览"),
        ("04", "核心能力", "规划到评价闭环"),
        ("05", "差异化", "不止聊天机器人"),
        ("06", "边界下一步", "诚实完成度"),
    ]
    for i, (n, t, d) in enumerate(items):
        col, row = i % 3, i // 3
        l = M + Inches(col * 4.15)
        top = Inches(1.7) + Inches(row * 2.35)
        rect(s, l, top, Inches(3.95), Inches(2.1), WHITE, 0.12)
        rect(s, l, top, Inches(3.95), Inches(0.1), TEAL if row == 0 else ACCENT, 0)
        txt(s, l + Inches(0.3), top + Inches(0.4), Inches(3.3), Inches(0.35), n, 18, True, TEAL)
        txt(s, l + Inches(0.3), top + Inches(0.9), Inches(3.3), Inches(0.4), t, 20, True, NAVY)
        txt(s, l + Inches(0.3), top + Inches(1.4), Inches(3.3), Inches(0.4), d, 13, False, MUTED)
    footer(s, 3)

    # ========== 4 Pain / Goal ==========
    s = prs.slides.add_slide(blank)
    bg(s, CREAM)
    title_block(s, "背景与目标", "从「能聊天」到「能出行」")
    rect(s, M, Inches(1.55), Inches(5.85), Inches(5.0), WHITE, 0.12)
    txt(s, Inches(1.05), Inches(1.85), Inches(5), Inches(0.4), "用户痛点", 16, True, ACCENT)
    pains = ["攻略分散，拼不出可执行行程", "大模型易编造店名，可信度低", "缺少地图与导航等落地动作", "没有反馈，下次依然盲目"]
    box = s.shapes.add_textbox(Inches(1.05), Inches(2.5), Inches(5.2), Inches(3.6))
    tf = box.text_frame
    tf.word_wrap = True
    for i, t in enumerate(pains):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(12)
        r = p.add_run()
        r.text = f"▸  {t}"
        font(r, 15, False, INK)

    rect(s, Inches(6.85), Inches(1.55), Inches(5.85), Inches(5.0), TEAL_DEEP, 0.12)
    txt(s, Inches(7.25), Inches(1.85), Inches(5), Inches(0.4), "项目目标", 16, True, TEAL)
    goals = ["可安装手机 App（Expo / APK）", "结构化多日行程，一次成型", "高德 POI + 和风天气回填", "评价闭环，越用越懂偏好"]
    box = s.shapes.add_textbox(Inches(7.25), Inches(2.5), Inches(5.2), Inches(3.6))
    tf = box.text_frame
    tf.word_wrap = True
    for i, t in enumerate(goals):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(12)
        r = p.add_run()
        r.text = f"▸  {t}"
        font(r, 15, False, WHITE)
    footer(s, 4)

    # ========== 5 Architecture ==========
    s = prs.slides.add_slide(blank)
    bg(s, CREAM)
    title_block(s, "产品方案", "三层分工：体验层 · 大脑层 · 能力层")
    layers = [
        ("手机 App", "规划 / 预览 / 地图\n社区 · 长辈模式", TEAL),
        ("Node 后端", "鉴权 · Agent · 校验\n限流 · Enrich", TEAL_DEEP),
        ("能力层", "Claude · 高德\n和风 · SQLite", ACCENT),
    ]
    for i, (title, desc, color) in enumerate(layers):
        l = M + Inches(i * 4.15)
        rect(s, l, Inches(1.7), Inches(3.95), Inches(2.7), color, 0.12)
        txt(s, l + Inches(0.35), Inches(2.05), Inches(3.3), Inches(0.45), title, 20, True, WHITE)
        box = s.shapes.add_textbox(l + Inches(0.35), Inches(2.7), Inches(3.3), Inches(1.3))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = desc
        font(r, 14, False, WHITE)
        if i < 2:
            txt(s, l + Inches(3.7), Inches(2.75), Inches(0.4), Inches(0.4), "→", 20, True, MUTED)
    rect(s, M, Inches(4.8), Inches(12.0), Inches(1.55), WHITE, 0.12)
    txt(s, Inches(1.0), Inches(5.05), Inches(11.5), Inches(0.35), "主流程", 13, True, TEAL_DARK)
    txt(
        s,
        Inches(1.0),
        Inches(5.45),
        Inches(11.5),
        Inches(0.6),
        "鉴权 → 填表 → Agent 调天气/POI → JSON 校验 → 回填坐标外链 → 预览保存 → 评价反哺",
        14,
        False,
        INK,
    )
    footer(s, 5)

    # ========== 6 UI Showcase hero ==========
    s = prs.slides.add_slide(blank)
    bg(s, TEAL_DEEP)
    txt(s, M, Inches(0.35), Inches(12), Inches(0.35), "界面演示", 13, False, TEAL)
    txt(s, M, Inches(0.7), Inches(12), Inches(0.5), "从引导到行程管理", 26, True, WHITE)
    labels = [("轻引导首页", "一句话说出想去哪"), ("我的行程", "云端同步与筛选"), ("社区发现", "行程 / 旅拍瀑布流")]
    phones = [hq_guide, hq_trips, hq_comm]
    for i, p in enumerate(phones):
        l = Inches(0.7 + i * 4.15)
        if p.exists():
            pic(s, p, l + Inches(0.35), Inches(1.35), height=Inches(4.75))
        t, d = labels[i]
        txt(s, l, Inches(6.25), Inches(3.9), Inches(0.3), t, 14, True, WHITE, PP_ALIGN.CENTER)
        txt(s, l, Inches(6.55), Inches(3.9), Inches(0.3), d, 12, False, RGBColor(0x9A, 0xC4, 0xC0), PP_ALIGN.CENTER)
    footer(s, 6, light=True)

    # ========== 7 More UI ==========
    s = prs.slides.add_slide(blank)
    bg(s, CREAM)
    title_block(s, "更多界面", "AI 聊聊 · 开机品牌 · 行程列表")
    more = [
        (hq_chat, "AI 聊聊", "流式对话，问签证美食与住宿"),
        (hq_splash if hq_splash.exists() else splash_phone, "开机页", "品牌第一印象"),
        (hq_trips, "行程库", "多城行程一览与管理"),
    ]
    for i, (p, title, desc) in enumerate(more):
        l = Inches(0.7 + i * 4.15)
        if p.exists():
            pic(s, p, l + Inches(0.45), Inches(1.4), height=Inches(4.7))
        txt(s, l, Inches(6.25), Inches(3.9), Inches(0.3), title, 15, True, NAVY, PP_ALIGN.CENTER)
        txt(s, l, Inches(6.55), Inches(3.9), Inches(0.3), desc, 12, False, MUTED, PP_ALIGN.CENTER)
    footer(s, 7)

    # ========== 8 Features ==========
    s = prs.slides.add_slide(blank)
    bg(s, CREAM)
    title_block(s, "核心能力", "主链路可演示，工程与体验兼顾")
    feats = [
        ("结构化行程", "目的地 / 天数 / 预算 / 偏好\n一次生成多日计划"),
        ("真实数据", "高德 POI + 和风天气\n坐标评分地址回填"),
        ("地图落地", "日计划时间轴\n导航 / 小红书 / 点评"),
        ("AI 闲聊", "SSE 流式问答\n支持拍照提问"),
        ("评价闭环", "赞踩标签反哺\n偏好注入下次生成"),
        ("行程管理", "本地 + 云端\n编辑分享导出"),
        ("长辈模式", "大字号 · 简表单\nTTS 朗读"),
        ("社区与画像", "Feed · 主动提醒\n旅行画像"),
    ]
    for i, (title, desc) in enumerate(feats):
        col, row = i % 4, i // 4
        l = M + Inches(col * 3.1)
        top = Inches(1.55) + Inches(row * 2.45)
        rect(s, l, top, Inches(2.95), Inches(2.2), WHITE, 0.12)
        rect(s, l, top, Inches(2.95), Inches(0.08), TEAL, 0)
        txt(s, l + Inches(0.22), top + Inches(0.35), Inches(2.5), Inches(0.4), title, 15, True, NAVY)
        box = s.shapes.add_textbox(l + Inches(0.22), top + Inches(0.9), Inches(2.5), Inches(1.0))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = desc
        font(r, 12, False, MUTED)
    footer(s, 8)

    # ========== 9 USP ==========
    s = prs.slides.add_slide(blank)
    bg(s, CREAM)
    title_block(s, "差异化卖点", "不是又一个旅游聊天机器人")
    usps = [
        ("01", "Agent + 真 POI", "工具检索带 poiId 回填，少编造"),
        ("02", "评价双闭环", "个人偏好 + 社区口碑共同选点"),
        ("03", "规划闲聊分离", "表单一次成型，结果可扫读"),
        ("04", "生成到落地", "地图与导航跳转，不只文字攻略"),
        ("05", "信任设计", "能力边界与隐私同意写进产品"),
        ("06", "工程完整", "Monorepo · 限流 · APK 可安装"),
    ]
    for i, (n, t, d) in enumerate(usps):
        top = Inches(1.5) + Inches(i * 0.82)
        rect(s, M, top, Inches(12.0), Inches(0.72), WHITE if i % 2 == 0 else CREAM_DEEP, 0.12)
        txt(s, Inches(0.95), top + Inches(0.2), Inches(0.7), Inches(0.35), n, 15, True, TEAL)
        txt(s, Inches(1.9), top + Inches(0.2), Inches(3.2), Inches(0.35), t, 15, True, NAVY)
        txt(s, Inches(5.3), top + Inches(0.2), Inches(7), Inches(0.35), d, 14, False, MUTED)
    footer(s, 9)

    # ========== 10 Demo script ==========
    s = prs.slides.add_slide(blank)
    bg(s, CREAM)
    title_block(s, "现场 Demo 剧本", "2–3 分钟走完主链路")
    steps = [
        ("1", "轻引导 / 规划", "目的地与偏好"),
        ("2", "生成预览", "天气 · 日计划 · 地图"),
        ("3", "落地跳转", "打开导航"),
        ("4", "保存同步", "行程库可见"),
        ("5", "评价反哺", "赞踩写入画像"),
        ("6", "再生成", "看闭环差异"),
    ]
    for i, (n, t, d) in enumerate(steps):
        top = Inches(1.5) + Inches(i * 0.8)
        rect(s, M, top, Inches(7.5), Inches(0.7), WHITE, 0.12)
        txt(s, Inches(0.95), top + Inches(0.18), Inches(0.5), Inches(0.35), n, 18, True, TEAL)
        txt(s, Inches(1.7), top + Inches(0.18), Inches(2.5), Inches(0.35), t, 15, True, NAVY)
        txt(s, Inches(4.3), top + Inches(0.18), Inches(3.5), Inches(0.35), d, 14, False, MUTED)
    if hq_guide.exists():
        pic(s, hq_guide, Inches(8.9), Inches(1.35), height=Inches(5.35))
    footer(s, 10)

    # ========== 11 Boundary ==========
    s = prs.slides.add_slide(blank)
    bg(s, CREAM)
    title_block(s, "完成度 · 边界 · 下一步", "主动说清边界")
    blocks = [
        (TEAL, "已完成", WHITE, ["Agent 行程主链路", "地图 / 外链 / 天气", "评价闭环与偏好", "闲聊 SSE + 拍照", "社区与长辈模式", "Android APK"]),
        (ACCENT, "明确不做", WHITE, ["不订票、不付款", "非自建导航引擎", "信息不保证 100%", "闭环需真实评价", "非生产级部署"]),
        (TEAL_DEEP, "下一步", TEAL, ["HTTPS 云部署", "云数据库", "可选 RAG", "PDF / 长图", "更多真用户数据"]),
    ]
    for i, (head, title, title_c, items) in enumerate(blocks):
        l = M + Inches(i * 4.15)
        rect(s, l, Inches(1.55), Inches(3.95), Inches(5.0), WHITE, 0.12)
        rect(s, l, Inches(1.55), Inches(3.95), Inches(0.6), head, 0.12)
        txt(s, l + Inches(0.3), Inches(1.68), Inches(3.3), Inches(0.4), title, 16, True, WHITE)
        box = s.shapes.add_textbox(l + Inches(0.35), Inches(2.45), Inches(3.3), Inches(3.8))
        tf = box.text_frame
        tf.word_wrap = True
        for j, t in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_before = Pt(10)
            r = p.add_run()
            r.text = f"·  {t}"
            font(r, 14, False, INK)
    footer(s, 11)

    # ========== 12 Closing ==========
    s = prs.slides.add_slide(blank)
    bg(s, TEAL_DEEP)
    rect(s, 0, 0, Inches(0.22), H, TEAL, 0)
    txt(s, M, Inches(1.8), Inches(7.5), Inches(0.4), "一句话收尾", 14, False, TEAL)
    txt(
        s,
        M,
        Inches(2.4),
        Inches(7.8),
        Inches(1.8),
        "途灵：能规划、能落地、\n能评价反哺的完整 Demo。",
        28,
        True,
        WHITE,
    )
    txt(
        s,
        M,
        Inches(4.5),
        Inches(7.5),
        Inches(0.8),
        "演示：引导生成 → 地图跳转 → 评价闭环\n边界：不做订票 · 非生产部署",
        14,
        False,
        RGBColor(0xA8, 0xC8, 0xC4),
    )
    txt(s, M, Inches(6.2), Inches(7.5), Inches(0.4), "Thank you · 欢迎提问与试用反馈", 16, False, WHITE)
    if hq_chat.exists():
        pic(s, hq_chat, Inches(9.1), Inches(1.1), height=Inches(5.5))

    out1 = DOCS / "途灵-项目汇报.pptx"
    out2 = DOCS / "tuling-project-pitch.pptx"
    prs.save(str(out1))
    prs.save(str(out2))
    print("Saved", out1)
    print("Saved", out2)


if __name__ == "__main__":
    build()
